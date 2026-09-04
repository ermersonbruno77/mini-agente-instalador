#!/usr/bin/env python3
"""Le um .pptx e devolve o texto de cada slide, em markdown simples.

Uso: python3 tools/pptx.py <arquivo.pptx> [--notas]

Escrito em 23/08/2026 para substituir o MarkItDown, que o Chefe mandou desinstalar.
Cobre o que o MarkItDown cobria: texto solto, caixa agrupada (recursivo), tabela e
os rotulos e valores de grafico, que e onde mora o numero dos decks dele.
"""
import sys
from pptx import Presentation
from pptx.util import Emu

def shapes_planas(shapes):
    for sh in shapes:
        if sh.shape_type == 6 and hasattr(sh, 'shapes'):   # GROUP
            yield from shapes_planas(sh.shapes)
        else:
            yield sh

def texto_do_shape(sh, buf):
    if sh.has_text_frame and sh.text_frame.text.strip():
        buf.append(sh.text_frame.text.strip())
    if getattr(sh, 'has_table', False) and sh.has_table:
        for r in sh.table.rows:
            buf.append('| ' + ' | '.join(c.text.strip().replace('\n', ' ') for c in r.cells) + ' |')
    if getattr(sh, 'has_chart', False) and sh.has_chart:
        ch = sh.chart
        try:
            cats = [str(c) for c in ch.plots[0].categories]
        except Exception:
            cats = []
        buf.append(f'[grafico: {ch.chart_type}]')
        if cats:
            buf.append('| serie | ' + ' | '.join(cats) + ' |')
        for s in ch.series:
            vals = ['' if v is None else str(v) for v in s.values]
            buf.append(f'| {s.name} | ' + ' | '.join(vals) + ' |')

def main():
    if len(sys.argv) < 2:
        print(__doc__); sys.exit(1)
    caminho = sys.argv[1]
    notas = '--notas' in sys.argv
    pres = Presentation(caminho)
    buf = []
    for i, slide in enumerate(pres.slides, 1):
        buf.append(f'\n## slide {i}')
        for sh in shapes_planas(slide.shapes):
            texto_do_shape(sh, buf)
        if notas and slide.has_notes_slide:
            t = slide.notes_slide.notes_text_frame.text.strip()
            if t:
                buf.append(f'> nota: {t}')
    print('\n'.join(buf))

if __name__ == '__main__':
    main()
